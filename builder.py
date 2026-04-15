"""PPTX builder for HRI Deck Builder. Converts SlideSpec objects into a branded PPTX."""
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from parser import SlideSpec


# Brand colors
TIDE = RGBColor(0x1E, 0x57, 0x73)
RISING_SUN = RGBColor(0xF2, 0x60, 0x44)
GOLDEN = RGBColor(0xBC, 0x89, 0x3F)
RESTORE = RGBColor(0x97, 0xCB, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# KPI card color rotation
KPI_COLORS = [TIDE, RISING_SUN, GOLDEN, RESTORE]

# Section color → Master 2 layout index
SECTION_COLOR_MAP = {
    "restore": 0,
    "golden": 1,
    "tide": 2,
    "risingsun": 3,
}

# Font name used in template
FONT_NAME = "Golos Text"


def build_deck(slides: list, template_path: str) -> bytes:
    """Build a PPTX from SlideSpec objects using the branded template.

    Args:
        slides: List of SlideSpec objects from the parser.
        template_path: Path to the presentation_template.pptx file.

    Returns:
        PPTX file contents as bytes.
    """
    prs = Presentation(template_path)
    original_count = len(prs.slides)

    for spec in slides:
        _add_slide(prs, spec)

    # Delete original sample slides (backwards to avoid index shifting)
    ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for i in range(original_count - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].get(f"{ns}id")
        if rId:
            prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_slide(prs: Presentation, spec: SlideSpec):
    """Add a single slide to the presentation based on the SlideSpec."""
    handler = {
        "TITLE": _add_title_slide,
        "SECTION": _add_section_slide,
        "CONTENT": _add_content_slide,
        "TWO-COL": _add_two_col_slide,
        "KPI": _add_kpi_slide,
        "QUOTE": _add_quote_slide,
        "PHOTO": _add_photo_slide,
        "BLANK": _add_blank_slide,
    }
    fn = handler.get(spec.slide_type, _add_content_slide)
    fn(prs, spec)


def _get_layout(prs: Presentation, master_idx: int, layout_idx: int):
    """Get a layout from a specific master by index."""
    return prs.slide_masters[master_idx].slide_layouts[layout_idx]


def _find_placeholder(slide, idx: int):
    """Find a placeholder by idx via iteration (dict-style access fails on Google Slides templates)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _replace_txBody(slide, ph, text_lines: list):
    """Replace a placeholder's txBody with the layout's, injecting text.

    python-pptx's .text setter creates a minimal txBody with bare <a:bodyPr/>
    and <a:lstStyle/>, which blocks inheritance of the layout's formatting
    (anchor, normAutofit, line spacing, font size, color). Instead, we deep-copy
    the layout's entire txBody (preserving all formatting) and inject text into it.

    Args:
        slide: The slide object (used to find the layout placeholder).
        ph: The slide-level placeholder shape to modify.
        text_lines: List of strings, one per paragraph.
    """
    from copy import deepcopy
    from lxml import etree

    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    pns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    ph_idx = ph.placeholder_format.idx
    layout_ph = _find_placeholder(slide.slide_layout, ph_idx)

    if layout_ph is None:
        ph.text = "\n".join(text_lines)
        return

    layout_txBody = layout_ph._sp.find(f"{pns}txBody")
    if layout_txBody is None:
        ph.text = "\n".join(text_lines)
        return

    # Deep-copy the layout's txBody (has correct bodyPr, lstStyle, paragraph fmt)
    new_txBody = deepcopy(layout_txBody)

    # Fix Google Slides line spacing artifacts in lstStyle.
    # Google Slides exports title placeholders with 222% line spacing which
    # creates huge gaps between wrapped lines in PowerPoint. Override to 100%.
    for lnSpc in new_txBody.iter(f"{ns}lnSpc"):
        spcPct = lnSpc.find(f"{ns}spcPct")
        if spcPct is not None:
            val = int(spcPct.get("val", "100000"))
            if val > 150000:  # anything over 150% is a Google Slides artifact
                spcPct.set("val", "100000")

    # Get the first paragraph as a formatting template
    existing_paras = new_txBody.findall(f"{ns}p")
    template_p = existing_paras[0] if existing_paras else None

    # Remove all existing paragraphs
    for p in existing_paras:
        new_txBody.remove(p)

    # Create one paragraph per text line, cloning formatting from template
    for line in text_lines:
        if template_p is not None:
            new_p = deepcopy(template_p)
            # Remove existing runs and end-para-rpr
            for r in new_p.findall(f"{ns}r"):
                new_p.remove(r)
        else:
            new_p = etree.Element(f"{ns}p")

        r_elem = etree.SubElement(new_p, f"{ns}r")
        t_elem = etree.SubElement(r_elem, f"{ns}t")
        t_elem.text = line
        new_txBody.append(new_p)

    # Replace the slide placeholder's txBody
    old_txBody = ph._sp.find(f"{pns}txBody")
    if old_txBody is not None:
        parent = ph._sp
        pos = list(parent).index(old_txBody)
        parent.remove(old_txBody)
        parent.insert(pos, new_txBody)


def _set_title(slide, text: str):
    """Set the title placeholder text, preserving layout formatting."""
    ph = _find_placeholder(slide, 0)
    if ph:
        _replace_txBody(slide, ph, [text])


def _set_bullets(slide, placeholder, bullets: list):
    """Set bullet points, preserving layout formatting."""
    _replace_txBody(slide, placeholder, bullets)


# --- Slide type handlers ---

def _add_title_slide(prs: Presentation, spec: SlideSpec):
    """TITLE slide: Master 0, Layout 0. Center title + subtitle."""
    layout = _get_layout(prs, 0, 0)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, spec.title)

    ph1 = _find_placeholder(slide, 1)
    if ph1:
        subtitle_parts = []
        if spec.metadata.get("subtitle"):
            subtitle_parts.append(spec.metadata["subtitle"])
        if spec.metadata.get("date"):
            subtitle_parts.append(spec.metadata["date"])
        if subtitle_parts:
            _replace_txBody(slide, ph1, [" | ".join(subtitle_parts)])


def _add_section_slide(prs: Presentation, spec: SlideSpec):
    """SECTION slide: Master 2, layout varies by color."""
    color = spec.metadata.get("color", "tide").lower().strip()
    layout_idx = SECTION_COLOR_MAP.get(color, 2)  # Default to Tide
    layout = _get_layout(prs, 2, layout_idx)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, spec.title)


def _add_content_slide(prs: Presentation, spec: SlideSpec):
    """CONTENT slide: Master 1, Layout 0 (OBJECT). Title + bullets."""
    layout = _get_layout(prs, 1, 0)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, spec.title)

    ph1 = _find_placeholder(slide, 1)
    if spec.bullets and ph1:
        _set_bullets(slide, ph1, spec.bullets)


def _add_two_col_slide(prs: Presentation, spec: SlideSpec):
    """TWO-COL slide: Master 1, Layout 15 (TWO_OBJECTS). Title + left/right bullets."""
    layout = _get_layout(prs, 1, 15)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, spec.title)

    left_bullets = spec.metadata.get("left_bullets", [])
    right_bullets = spec.metadata.get("right_bullets", [])

    ph1 = _find_placeholder(slide, 1)
    ph2 = _find_placeholder(slide, 2)
    if left_bullets and ph1:
        _set_bullets(slide, ph1, left_bullets)
    if right_bullets and ph2:
        _set_bullets(slide, ph2, right_bullets)


def _add_kpi_slide(prs: Presentation, spec: SlideSpec):
    """KPI slide: Master 1, Layout 3 (TITLE_ONLY). Title + programmatic metric cards."""
    layout = _get_layout(prs, 1, 3)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, spec.title)

    if not spec.kpi_items:
        return

    # Card layout constants
    num_cards = len(spec.kpi_items)
    card_width = Inches(2.8)
    card_height = Inches(2.5)
    spacing = Inches(0.3)
    total_width = num_cards * card_width + (num_cards - 1) * spacing
    start_x = (prs.slide_width - total_width) // 2
    top = Inches(2.5)
    corner_radius = Inches(0.15)

    for i, (value, label) in enumerate(spec.kpi_items):
        color = KPI_COLORS[i % len(KPI_COLORS)]
        left = start_x + i * (card_width + spacing)

        # Rounded rectangle card
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, card_width, card_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        # Set corner radius
        shape.adjustments[0] = 0.05

        # Add text to the shape
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        # Value paragraph (large, bold)
        p_val = tf.paragraphs[0]
        p_val.alignment = PP_ALIGN.CENTER
        run_val = p_val.add_run()
        run_val.text = value
        run_val.font.size = Pt(36)
        run_val.font.bold = True
        run_val.font.color.rgb = WHITE
        run_val.font.name = FONT_NAME

        # Label paragraph (smaller)
        p_label = tf.add_paragraph()
        p_label.alignment = PP_ALIGN.CENTER
        run_label = p_label.add_run()
        run_label.text = label
        run_label.font.size = Pt(14)
        run_label.font.bold = False
        run_label.font.color.rgb = WHITE
        run_label.font.name = FONT_NAME

        # Vertical centering
        tf.paragraphs[0].space_before = Pt(30)


def _add_quote_slide(prs: Presentation, spec: SlideSpec):
    """QUOTE slide: Master 1, Layout 2 (Pullquote). Title + quote text + attribution."""
    layout = _get_layout(prs, 1, 2)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, spec.title)

    quote_text = spec.metadata.get("quote", "")
    attribution = spec.metadata.get("attribution", "")

    ph1 = _find_placeholder(slide, 1)
    if ph1:
        lines = []
        if quote_text:
            lines.append(f"\u201c{quote_text}\u201d")
        if attribution:
            lines.append(f"\u2014 {attribution}")
        if lines:
            _replace_txBody(slide, ph1, lines)


def _add_photo_slide(prs: Presentation, spec: SlideSpec):
    """PHOTO slide: Master 1, Layout 4 (Photo 1). Title + bullets (photo placeholder left empty)."""
    layout = _get_layout(prs, 1, 4)
    slide = prs.slides.add_slide(layout)
    _set_title(slide, spec.title)

    ph1 = _find_placeholder(slide, 1)
    if spec.bullets and ph1:
        _set_bullets(slide, ph1, spec.bullets)


def _add_blank_slide(prs: Presentation, spec: SlideSpec):
    """BLANK slide: Master 1, Layout 17. No content."""
    layout = _get_layout(prs, 1, 17)
    prs.slides.add_slide(layout)
