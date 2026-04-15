"""Tests for the PPTX builder."""
import sys
import os
import io
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from parser import parse_markdown
from builder import build_deck

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                              "assets", "presentation_template.pptx")
SAMPLE_PATH = os.path.join(os.path.dirname(__file__), "fixtures", "sample_input.md")


def _load_sample():
    with open(SAMPLE_PATH) as f:
        md = f.read()
    return parse_markdown(md)


def _build_and_load(slides):
    pptx_bytes = build_deck(slides, TEMPLATE_PATH)
    return Presentation(io.BytesIO(pptx_bytes))


def _find_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def test_output_opens_as_valid_pptx():
    slides = _load_sample()
    pptx_bytes = build_deck(slides, TEMPLATE_PATH)
    # Should not raise
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert prs is not None


def test_slide_count_matches_input():
    slides = _load_sample()
    prs = _build_and_load(slides)
    assert len(prs.slides) == len(slides)


def test_slide_dimensions():
    slides = _load_sample()
    prs = _build_and_load(slides)
    # 13.333" x 7.500"
    assert abs(prs.slide_width / 914400 - 13.333) < 0.01
    assert abs(prs.slide_height / 914400 - 7.500) < 0.01


def test_no_sample_slides_remain():
    slides = _load_sample()
    prs = _build_and_load(slides)
    # Original template has 12 sample slides; output should have exactly len(slides)
    assert len(prs.slides) == 10


def test_title_text_on_slides():
    slides = _load_sample()
    prs = _build_and_load(slides)
    for i, (spec, slide) in enumerate(zip(slides, prs.slides)):
        if spec.slide_type == "BLANK":
            continue
        ph = _find_ph(slide, 0)
        assert ph is not None, f"Slide {i} ({spec.slide_type}) has no title placeholder"
        assert ph.text == spec.title, f"Slide {i}: expected '{spec.title}', got '{ph.text}'"


def test_content_slide_has_bullets():
    slides = _load_sample()
    prs = _build_and_load(slides)
    # Slide 2 is CONTENT with 3 bullets
    slide = prs.slides[2]
    ph = _find_ph(slide, 1)
    assert ph is not None
    paragraphs = [p.text for p in ph.text_frame.paragraphs if p.text.strip()]
    assert len(paragraphs) == 3


def test_kpi_slide_has_shapes():
    slides = _load_sample()
    prs = _build_and_load(slides)
    # Slide 3 is KPI with 4 items
    slide = prs.slides[3]
    # Count non-placeholder shapes with text (the KPI cards)
    kpi_shapes = []
    for s in slide.shapes:
        if s.has_text_frame:
            try:
                s.placeholder_format
            except ValueError:
                if s.text_frame.text.strip():
                    kpi_shapes.append(s)
    assert len(kpi_shapes) == 4


def test_two_col_slide_has_both_columns():
    slides = _load_sample()
    prs = _build_and_load(slides)
    # Slide 4 is TWO-COL
    slide = prs.slides[4]
    ph1 = _find_ph(slide, 1)
    ph2 = _find_ph(slide, 2)
    assert ph1 is not None, "Left column placeholder missing"
    assert ph2 is not None, "Right column placeholder missing"
    left_paras = [p.text for p in ph1.text_frame.paragraphs if p.text.strip()]
    right_paras = [p.text for p in ph2.text_frame.paragraphs if p.text.strip()]
    assert len(left_paras) == 3
    assert len(right_paras) == 3


def test_quote_slide_has_text():
    slides = _load_sample()
    prs = _build_and_load(slides)
    # Slide 5 is QUOTE
    slide = prs.slides[5]
    ph = _find_ph(slide, 1)
    assert ph is not None
    text = ph.text_frame.text
    assert "failed to execute" in text
    assert "Bill Simmons" in text


def test_section_slides_use_correct_layouts():
    slides = _load_sample()
    prs = _build_and_load(slides)
    # Slide 1 is SECTION (tide), slide 6 is SECTION (golden)
    assert prs.slides[1].slide_layout.name == "Tide"
    assert prs.slides[6].slide_layout.name == "Golden"


def test_layout_assignments():
    slides = _load_sample()
    prs = _build_and_load(slides)
    expected_layouts = [
        "TITLE",         # 0: TITLE from frontmatter
        "Tide",          # 1: SECTION tide
        "OBJECT",        # 2: CONTENT
        "TITLE_ONLY",    # 3: KPI
        "TWO_OBJECTS",   # 4: TWO-COL
        "Pullquote",     # 5: QUOTE
        "Golden",        # 6: SECTION golden
        "OBJECT",        # 7: CONTENT
        "OBJECT",        # 8: CONTENT
        "BLANK",         # 9: BLANK
    ]
    for i, (slide, expected) in enumerate(zip(prs.slides, expected_layouts)):
        assert slide.slide_layout.name == expected, \
            f"Slide {i}: expected layout '{expected}', got '{slide.slide_layout.name}'"


def test_empty_input_produces_empty_deck():
    slides = parse_markdown("")
    pptx_bytes = build_deck(slides, TEMPLATE_PATH)
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) == 0
